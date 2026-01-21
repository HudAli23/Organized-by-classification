import logging
from transformers import pipeline
import torch
import fitz  # PyMuPDF for PDF extraction
import docx  # python-docx for DOCX extraction
import pptx  # python-pptx for PPTX extraction
import openpyxl  # for Excel files
from bs4 import BeautifulSoup  # for HTML files
import chardet  # for detecting file encodings
import re
import os
import json
from concurrent.futures import ThreadPoolExecutor, as_completed
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
import networkx as nx
from collections import defaultdict
import numpy as np
from datetime import datetime
import hashlib
import threading
import csv
import shutil
import nltk

# Shared constants
# ... existing code ...

# Domain-specific categorization hints
DOMAIN_INDICATORS = {
    "Cybersecurity": {
        "patterns": [
            "security", "vulnerability", "threat", "exploit", "malware",
            "firewall", "encryption", "cyber", "attack", "defense",
            "incident", "response", "penetration", "testing", "audit"
        ],
        "priority": 10  # Higher priority for security-related content
    },
    "Development": {
        "patterns": [
            "development", "code", "programming", "software", "application",
            "frontend", "backend", "api", "database", "interface"
        ],
        "priority": 8
    },
    "Data Science": {
        "patterns": [
            "data science", "machine learning", "analytics", "statistics",
            "model", "prediction", "analysis", "dataset", "training"
        ],
        "priority": 8
    },
    "DevOps": {
        "patterns": [
            "deployment", "pipeline", "infrastructure", "container",
            "kubernetes", "docker", "ci/cd", "automation", "monitoring"
        ],
        "priority": 8
    }
}

# Default classification labels
LABELS = [
    "Documents",
    "Spreadsheets",
    "Presentations",
    "Images",
    "Code",
    "Audio",
    "Video",
    "Archives",
    "Research",
    "Reports",
    "Data",
    "Documentation",
    "Academic",
    "Others"
]

def initialize_components():
    """Initialize all necessary components."""
    global document_classifier_gpu, document_classifier_cpu
    
    # Initialize document classifier
    if document_classifier_gpu is None or document_classifier_cpu is None:
        try:
            # Initialize GPU classifier if available
            if torch.cuda.is_available():
                device = get_optimal_device()
                logging.info(f"Initializing GPU classifier on device {device}")
                document_classifier_gpu = pipeline(
                    "zero-shot-classification",
                    model="facebook/bart-large-mnli",
                    device=device,
                    torch_dtype=torch.float16  # Use FP16 for GPU
                )
                logging.info("GPU classifier initialized successfully")
            
            # Initialize CPU classifier
            logging.info("Initializing CPU classifier")
            document_classifier_cpu = pipeline(
                "zero-shot-classification",
                model="facebook/bart-large-mnli",
                device=-1,  # Force CPU
                torch_dtype=torch.float32
            )
            logging.info("CPU classifier initialized successfully")
            return True
        except Exception as e:
            logging.error(f"Error initializing document classifier: {e}")
            return False
    return True

# Initialize the document classifier
document_classifier_gpu = None
document_classifier_cpu = None
DEVICE = None

def get_optimal_device():
    """Determine the optimal device (CUDA GPU if available, otherwise CPU)."""
    try:
        if torch.cuda.is_available():
            # Get the GPU with most free memory
            device_count = torch.cuda.device_count()
            if device_count > 0:
                free_memory = []
                for i in range(device_count):
                    torch.cuda.set_device(i)
                    torch.cuda.empty_cache()
                    free_memory.append(torch.cuda.get_device_properties(i).total_memory - \
                                    torch.cuda.memory_allocated(i))
                best_gpu = free_memory.index(max(free_memory))
                logging.info(f"Using GPU {best_gpu} with {free_memory[best_gpu]/1024**2:.1f}MB free memory")
                return best_gpu
        logging.info("No GPU available, using CPU")
        return -1
    except Exception as e:
        logging.error(f"Error detecting device: {e}")
        return -1

def initialize_document_classifier():
    """Initialize the document classifier with proper error handling."""
    global document_classifier_gpu, document_classifier_cpu, DEVICE
    
    if document_classifier_gpu is not None and document_classifier_cpu is not None:
        return True
        
    try:
        # Initialize GPU classifier if available
        if torch.cuda.is_available():
            DEVICE = get_optimal_device()
            logging.info(f"Initializing GPU classifier on device {DEVICE}")
            document_classifier_gpu = pipeline(
                "zero-shot-classification",
                model="facebook/bart-large-mnli",
                device=DEVICE,
                torch_dtype=torch.float16  # Use FP16 for GPU
            )
            logging.info("GPU classifier initialized successfully")
        
        # Initialize CPU classifier
        logging.info("Initializing CPU classifier")
        document_classifier_cpu = pipeline(
            "zero-shot-classification",
            model="facebook/bart-large-mnli",
            device=-1,  # Force CPU
            torch_dtype=torch.float32
        )
        logging.info("CPU classifier initialized successfully")
        
        return True
    except Exception as e:
        logging.error(f"Error initializing document classifier: {e}")
        return False

# Initialize the classifier when the module is imported
initialize_document_classifier()

# Shared functions
def clean_text(text):
    """Clean and normalize text content with improved preprocessing."""
    if not text:
        return ""
    
    # Convert to lowercase for better matching
    text = text.lower()
    
    # Remove extra whitespace
    text = re.sub(r'\s+', ' ', text)
    
    # Keep meaningful punctuation and special characters
    text = re.sub(r'[^\w\s.,!?@#$%&*()-]', ' ', text)
    
    # Remove very short words (likely noise)
    words = text.split()
    words = [w for w in words if len(w) > 1]
    
    # Remove common file extensions that might confuse classification
    words = [w for w in words if not w.startswith('.') and not w.endswith(('.txt', '.pdf', '.doc', '.xls'))]
    
    # Join words back together
    text = ' '.join(words)
    
    return text.strip()

def extract_from_txt(path):
    """Extract text from a .txt file with improved encoding detection."""
    try:
        # First try to detect encoding
        with open(path, 'rb') as file:
            raw = file.read()
            if not raw:
                return ""
                
            result = chardet.detect(raw)
            encodings = [
                result['encoding'] if result and result['encoding'] else 'utf-8',
                'utf-8', 'ascii', 'iso-8859-1', 'cp1252'
            ]
            
            # Try each encoding
            for encoding in encodings:
                try:
                    if encoding:
                        text = raw.decode(encoding)
                        return clean_text(text)
                except Exception:
                    continue
            
            logging.warning(f"Could not decode text file {path} with any encoding")
            return ""
            
    except Exception as e:
        logging.warning(f"Error reading text file {path}: {e}")
        return ""

def extract_from_pdf(path):
    """Extract text from a .pdf file."""
    try:
        with fitz.open(path) as pdf:
            text = []
            for page in pdf:
                text.append(page.get_text())
            return clean_text("\n".join(text))
    except Exception as e:
        logging.warning(f"Error reading PDF file {path}: {e}")
        return ""

def extract_from_docx(path):
    """Extract text from a .docx file."""
    try:
        doc = docx.Document(path)
        text = []
        # Extract text from paragraphs
        for para in doc.paragraphs:
            text.append(para.text)
        # Extract text from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text.append(cell.text)
        return clean_text("\n".join(text))
    except Exception as e:
        logging.warning(f"Error reading Word file {path}: {e}")
        return ""

def extract_from_pptx(path):
    """Extract text from a .pptx file."""
    try:
        prs = pptx.Presentation(path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return clean_text("\n".join(text))
    except Exception as e:
        logging.warning(f"Error reading PowerPoint file {path}: {e}")
        return ""

def extract_from_excel(path):
    """Extract text from Excel files with timeout and improved handling."""
    try:
        import threading
        import queue
        import pandas as pd

        text = []
        result_queue = queue.Queue()

        def excel_worker():
            try:
                # Try with openpyxl first
                wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
                
                # Add file metadata
                text.append(f"Excel file: {os.path.basename(path)}")
                text.append(f"Number of sheets: {len(wb.sheetnames)}")
                
                for sheet in wb.sheetnames:
                    ws = wb[sheet]
                    text.append(f"\nSheet: {sheet}")
                    
                    # Get column headers
                    headers = []
                    for cell in next(ws.iter_rows()):
                        if cell.value:
                            headers.append(str(cell.value))
                    if headers:
                        text.append("Headers: " + " | ".join(headers))
                    
                    # Process rows
                    row_count = 0
                    for row in ws.iter_rows():
                        row_text = []
                        for cell in row:
                            if cell.value:
                                row_text.append(str(cell.value))
                        if row_text:
                            text.append(" | ".join(row_text))
                            row_count += 1
                    text.append(f"Total rows in sheet {sheet}: {row_count}")
                
                wb.close()
                result_queue.put((True, text))
            except Exception as e:
                logging.warning(f"openpyxl failed for {path}, trying pandas: {str(e)}")
                try:
                    # Read Excel file with pandas
                    all_sheets = pd.read_excel(path, sheet_name=None)
                    
                    # Add file metadata
                    text.append(f"Excel file: {os.path.basename(path)}")
                    text.append(f"Number of sheets: {len(all_sheets)}")
                    
                    for sheet_name, df in all_sheets.items():
                        text.append(f"\nSheet: {sheet_name}")
                        
                        # Get column headers
                        headers = df.columns.tolist()
                        if headers:
                            text.append("Headers: " + " | ".join(map(str, headers)))
                        
                        # Process rows
                        for _, row in df.iterrows():
                            row_text = " | ".join(row.astype(str).values)
                            if row_text.strip():
                                text.append(row_text)
                        text.append(f"Total rows in sheet {sheet_name}: {len(df)}")
                    
                    result_queue.put((True, text))
                except Exception as e2:
                    logging.warning(f"pandas failed for {path}: {str(e2)}")
                    result_queue.put((False, str(e2)))

        # Start Excel processing in a separate thread
        excel_thread = threading.Thread(target=excel_worker)
        excel_thread.daemon = True
        excel_thread.start()

        # Wait for result with timeout (15 seconds)
        excel_thread.join(timeout=15)

        if excel_thread.is_alive():
            logging.warning(f"Excel processing timed out for {path}")
            return f"Excel file {os.path.basename(path)} (processing timed out)"

        # Get the result
        try:
            success, result = result_queue.get_nowait()
            if success:
                return "\n".join(result)
            else:
                logging.warning(f"Excel processing failed: {result}")
                return f"Excel file {os.path.basename(path)}"
        except queue.Empty:
            logging.warning(f"Excel processing result not available for {path}")
            return f"Excel file {os.path.basename(path)}"

    except Exception as e:
        logging.error(f"Failed to process Excel file {path}: {str(e)}")
        return f"Excel file {os.path.basename(path)}"

def extract_from_html(path):
    """Extract text from HTML files with improved encoding detection."""
    try:
        # First try to detect encoding
        with open(path, 'rb') as file:
            raw = file.read()
            if not raw:
                return ""
                
            result = chardet.detect(raw)
            encodings = [
                result['encoding'] if result and result['encoding'] else 'utf-8',
                'utf-8', 'ascii', 'iso-8859-1', 'cp1252'
            ]
            
            # Try each encoding
            for encoding in encodings:
                try:
                    if encoding:
                        content = raw.decode(encoding)
                        soup = BeautifulSoup(content, 'html.parser')
                        # Remove script and style elements
                        for script in soup(["script", "style"]):
                            script.decompose()
                        # Get text content
                        text = soup.get_text()
                        return clean_text(text)
                except Exception:
                    continue
            
            logging.warning(f"Could not decode HTML file {path} with any encoding")
            return ""
            
    except Exception as e:
        logging.warning(f"Error reading HTML file {path}: {e}")
        return ""

def extract_from_code(path):
    """Extract text from code files with improved encoding detection."""
    try:
        # First try to detect encoding
        with open(path, 'rb') as file:
            raw = file.read()
            if not raw:
                return ""
                
            result = chardet.detect(raw)
            encodings = [
                result['encoding'] if result and result['encoding'] else 'utf-8',
                'utf-8', 'ascii', 'iso-8859-1', 'cp1252'
            ]
            
            # Try each encoding
            for encoding in encodings:
                try:
                    if encoding:
                        content = raw.decode(encoding)
                        # Remove common code syntax but keep meaningful content
                        content = re.sub(r'//.*$', ' ', content, flags=re.MULTILINE)  # Remove single-line comments
                        content = re.sub(r'/\*.*?\*/', ' ', content, flags=re.DOTALL)  # Remove multi-line comments
                        content = re.sub(r'[{}()\[\];]', ' ', content)  # Remove common programming punctuation
                        return clean_text(content)
                except Exception:
                    continue
            
            logging.warning(f"Could not decode code file {path} with any encoding")
            return ""
            
    except Exception as e:
        logging.warning(f"Error reading code file {path}: {e}")
        return ""

def extract_text_from_file(path):
    """Extract text from various types of files with improved encoding detection."""
    if not path:
        return ""
        
    try:
        ext = os.path.splitext(path)[1].lower()
        
        # Skip binary and media files
        binary_extensions = {
            '.mp3', '.mp4', '.wav', '.avi', '.mov', '.flv', '.wmv',  # Audio/Video
            '.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp',  # Images
            '.zip', '.rar', '.7z', '.tar', '.gz', '.exe', '.dll',  # Binary/Archives
            '.db', '.dat', '.bin', '.pyc', '.pyo', '.pyd'  # Data files
        }
        
        if ext in binary_extensions or os.path.getsize(path) == 0:
            logging.info(f"Skipping binary or empty file: {path}")
            return ""
        
        # Dictionary mapping file extensions to their extraction functions
        extractors = {
            '.txt': extract_from_txt,
            '.pdf': extract_from_pdf,
            '.docx': extract_from_docx,
            '.doc': extract_from_docx,
            '.pptx': extract_from_pptx,
            '.ppt': extract_from_pptx,
            '.xlsx': extract_from_excel,
            '.xls': extract_from_excel,
            '.csv': extract_from_csv,
            '.html': extract_from_html,
            '.htm': extract_from_html,
            '.py': extract_from_code,
            '.java': extract_from_code,
            '.js': extract_from_code,
            '.cpp': extract_from_code,
            '.h': extract_from_code,
            '.css': extract_from_code,
            '.xml': extract_from_code,
            '.json': extract_from_code,
            '.md': extract_from_txt,
            '.rtf': extract_from_txt,
            '.php': extract_from_code,
            '.c': extract_from_code,
            '.cs': extract_from_code,
            '.rb': extract_from_code,
            '.swift': extract_from_code,
            '.m': extract_from_code,
            '.sql': extract_from_code,
            '.yaml': extract_from_code,
            '.yml': extract_from_code,
            '.ini': extract_from_txt,
            '.cfg': extract_from_txt,
            '.log': extract_from_txt,
            '.tex': extract_from_txt
        }
        
        try:
            if ext in extractors:
                logging.info(f"Extracting text from {path} using {extractors[ext].__name__}")
                text = extractors[ext](path)
                if text:
                    # Normalize whitespace and remove very short content
                    text = ' '.join(text.split())
                    if len(text) > 10:  # Only return if we have meaningful content
                        logging.info(f"Successfully extracted {len(text)} characters from {path}")
                        return text
                    else:
                        logging.warning(f"Extracted text too short from {path}: {text}")
                else:
                    logging.warning(f"No text extracted from {path}")
            else:
                # For unknown file types, try multiple encodings
                return extract_from_unknown(path)
                
        except Exception as e:
            logging.error(f"Error extracting text from {path}: {e}")
            return ""
            
    except Exception as e:
        logging.error(f"Error accessing file {path}: {e}")
        return ""

def extract_from_unknown(path):
    """Extract text from unknown file types with multiple encoding attempts."""
    try:
        # First try to detect encoding
        with open(path, 'rb') as f:
            raw_data = f.read()
            if not raw_data:
                return ""
                
            # Try chardet first
            result = chardet.detect(raw_data)
            encodings = [
                result['encoding'] if result and result['encoding'] else 'utf-8',
                'utf-8', 'ascii', 'iso-8859-1', 'cp1252', 'utf-16', 'utf-32'
            ]
            
            # Try each encoding
            for encoding in encodings:
                try:
                    if encoding:
                        text = raw_data.decode(encoding)
                        # Check if the decoded text looks valid
                        if any(ord(c) < 128 for c in text):
                            # Clean and normalize text
                            text = clean_text(text)
                            if len(text) > 10:
                                logging.info(f"Successfully extracted text from unknown file type using {encoding} encoding: {path}")
                                return text
                except Exception:
                    continue
            
            logging.warning(f"Could not extract text from unknown file type {path} - no valid encoding found")
            return ""
            
    except Exception as e:
        logging.error(f"Error extracting text from unknown file: {path} - {str(e)}")
        return ""

def extract_from_csv(path):
    """Extract text from CSV files with improved encoding detection."""
    try:
        # First try to detect encoding
        with open(path, 'rb') as file:
            raw = file.read()
            if not raw:
                return ""
                
            result = chardet.detect(raw)
            encodings = [
                result['encoding'] if result and result['encoding'] else 'utf-8',
                'utf-8', 'ascii', 'iso-8859-1', 'cp1252'
            ]
            
            # Try each encoding
            for encoding in encodings:
                try:
                    if encoding:
                        with open(path, 'r', encoding=encoding) as file:
                            reader = csv.reader(file)
                            rows = []
                            for row in reader:
                                rows.extend(str(cell) for cell in row if cell)
                            return ' '.join(rows)
                except Exception:
                    continue
            
            logging.warning(f"Could not decode CSV file {path} with any encoding")
            return ""
            
    except Exception as e:
        logging.warning(f"Error reading CSV file {path}: {e}")
        return ""

def classify_document(text, labels, threshold=0.3, chunk_size=512):
    """Classify the document into one of the predefined labels based on its content."""
    global document_classifier_gpu, document_classifier_cpu
    
    # Try to initialize classifiers if not already initialized
    if document_classifier_gpu is None or document_classifier_cpu is None:
        if not initialize_document_classifier():
            return "Others", 0.3
            
    try:
        # Log the start of classification
        logging.info("Starting document classification")
        # Clean and preprocess text
        text = text.strip()
        if not text:
            logging.warning("No text provided for classification.")
            return "Others", 0.3
            
        # Check for Excel/tabular data indicators first
        excel_indicators = [
            'excel file:', 'sheet:', 'headers:', 'total rows',
            'rank', 'score', 'total', 'year', 'column', 'row', 'sheet', 'table',
            'data', 'value', 'average', 'sum', 'count', 'percentage',
            '|',
            '.xlsx', '.xls'
        ]
        
        # Count Excel indicators
        text_lower = text.lower()
        indicator_count = sum(1 for indicator in excel_indicators if indicator in text_lower)
        logging.info(f"Found {indicator_count} Excel indicators in text")
        if indicator_count >= 2 or '.xls' in text_lower:
            logging.info(f"Detected Excel file based on {indicator_count} indicators")
            return "Spreadsheet", 0.9
            
        MAX_TEXT_LENGTH = 5000
        if len(text) > MAX_TEXT_LENGTH:
            logging.info(f"Text too long ({len(text)} chars), truncating to {MAX_TEXT_LENGTH} chars")
            first_part = text[:MAX_TEXT_LENGTH//2]
            last_part = text[-MAX_TEXT_LENGTH//2:]
            text = first_part + " ... " + last_part
            
        words = text.split()
        if not words:
            logging.warning("No words found in the text for classification.")
            return "Others", 0.3
            
        chunk_size = 128
        MAX_CHUNKS = 3
        chunks = []
        if len(words) > chunk_size:
            chunks.append(' '.join(words[:chunk_size]))
        if len(words) > chunk_size * 2:
            mid_start = len(words) // 2 - chunk_size // 2
            chunks.append(' '.join(words[mid_start:mid_start + chunk_size]))
        if len(words) > chunk_size:
            chunks.append(' '.join(words[-chunk_size:]))
        else:
            chunks = [' '.join(words)]
            
        logging.info(f"Processing {len(chunks)} text chunks")
        scores = {label: 0.0 for label in labels}
        processed_chunks = 0
        
        # Create a thread pool for parallel processing
        with ThreadPoolExecutor(max_workers=2) as executor:
            futures = []
            
            # Submit chunks to both GPU and CPU classifiers
            for chunk in chunks:
                if document_classifier_gpu is not None:
                    futures.append(executor.submit(document_classifier_gpu, chunk, labels))
                if document_classifier_cpu is not None:
                    futures.append(executor.submit(document_classifier_cpu, chunk, labels))
            
            # Process results as they complete
            for future in futures:
                try:
                    result = future.result(timeout=10)
                    processed_chunks += 1
                    for label, score in zip(result['labels'], result['scores']):
                        scores[label] += score
                except Exception as e:
                    logging.error(f"Error processing chunk: {e}")
                    continue
        
        if processed_chunks == 0:
            logging.warning("No chunks processed successfully, trying fallback classification")
            text_lower = text.lower()
            keyword_categories = {
                "Spreadsheet": ["excel", "spreadsheet", "table", "column", "row", "data", "sheet", ".xls"],
                "Research": ["research", "study", "analysis", "methodology", "findings"],
                "Report": ["report", "summary", "overview", "results", "statistics"],
                "Data": ["data", "statistics", "numbers", "metrics", "measurements"],
                "Documentation": ["documentation", "guide", "manual", "instructions"],
                "Academic": ["university", "academic", "education", "student", "faculty"]
            }
            category_scores = {cat: sum(1 for kw in kws if kw in text_lower) 
                             for cat, kws in keyword_categories.items()}
            if any(category_scores.values()):
                best_category = max(category_scores.items(), key=lambda x: x[1])[0]
                logging.info(f"Fallback classification found category: {best_category}")
                return best_category, 0.6
            logging.error("Fallback classification failed")
            return "Others", 0.3
            
        # Average the scores across all processed chunks
        for label in scores:
            scores[label] /= processed_chunks
            
        best_label = max(scores.items(), key=lambda x: x[1])[0]
        best_score = scores[best_label]
        
        if '.xls' in text_lower and best_score < 0.7:
            logging.info("Overriding classification to Spreadsheet based on file extension")
            return "Spreadsheet", 0.9
            
        logging.info(f"Classification completed. Best label: {best_label}, Score: {best_score}")
        if best_score >= threshold:
            return best_label, best_score
        else:
            logging.info(f"Score {best_score} below threshold {threshold}, returning Others")
            return "Others", 0.3
            
    except Exception as e:
        logging.error(f"Classification failed: {e}")
        return "Others", 0.3

def is_tesseract_in_path():
    """Check if Tesseract is available in system PATH."""
    try:
        import subprocess
        subprocess.run(['tesseract', '--version'], capture_output=True)
        return True
    except FileNotFoundError:
        return False

def configure_tesseract():
    """Configure Tesseract path and check installation."""
    import sys
    import os
    import pytesseract
    
    # Set Tesseract path directly
    pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
    
    try:
        # Verify Tesseract is working
        version = pytesseract.get_tesseract_version()
        logging.info(f"Tesseract OCR version {version} is properly configured.")
        return True
    except Exception as e:
        logging.error(f"Error verifying Tesseract: {e}")
        return False

def check_tesseract_installation():
    """Check if Tesseract is properly installed and provide guidance if it's not."""
    import sys
    import subprocess
    import webbrowser
    
    # First try to configure Tesseract
    if configure_tesseract():
        return True
        
    # If configuration failed, provide installation guidance
    logging.info("To install Tesseract OCR on Windows:")
    logging.info("1. Download Tesseract installer from: https://github.com/UB-Mannheim/tesseract/wiki")
    logging.info("2. Run the installer (make sure to check 'Add to system PATH' during installation)")
    logging.info("3. Restart your application after installation")
    
    # Ask if user wants to open the download page
    if input("Would you like to open the Tesseract download page? (y/n): ").lower() == 'y':
        webbrowser.open('https://github.com/UB-Mannheim/tesseract/wiki')
    
    return False

class ContentAnalyzer:
    def __init__(self):
        self.vectorizer = None  # Will be initialized on first use
        self.content_graph = nx.Graph()
        self.file_metadata = {}
        self.content_vectors = {}
        self.relationship_cache = {}
        self.feature_names = None
        self.is_fitted = False
        self.lock = threading.Lock()  # Add thread safety
        
        # Initialize GPU and CPU models
        self.gpu_model = None
        self.cpu_model = None
        self.initialize_models()

    def initialize_models(self):
        """Initialize both GPU and CPU models for parallel processing."""
        try:
            # Initialize GPU model if available
            if torch.cuda.is_available():
                self.gpu_model = pipeline(
                    "zero-shot-classification",
                    model="facebook/bart-large-mnli",
                    device=0,  # Use first GPU
                    torch_dtype=torch.float16  # Use FP16 for better performance
                )
                logging.info("GPU model initialized successfully")
            
            # Initialize CPU model
            self.cpu_model = pipeline(
                "zero-shot-classification",
                model="facebook/bart-large-mnli",
                device=-1,  # Force CPU
                torch_dtype=torch.float32
            )
            logging.info("CPU model initialized successfully")
            
        except Exception as e:
            logging.error(f"Error initializing models: {e}")
            
    def _initialize_vectorizer(self):
        """Initialize the vectorizer if not already done."""
        if self.vectorizer is None:
            self.vectorizer = TfidfVectorizer(
                max_features=5000,
                stop_words='english',
                ngram_range=(1, 3),
                lowercase=True,
                strip_accents='unicode'
            )

    def _extract_topics(self, text):
        """Extract main topics from text using TF-IDF."""
        try:
            # Clean and preprocess text
            text = clean_text(text)
            if not text:
                return []

            with self.lock:
                self._initialize_vectorizer()
                
                # If not fitted, fit the vectorizer
                if not self.is_fitted:
                    tfidf_matrix = self.vectorizer.fit_transform([text])
                    self.feature_names = self.vectorizer.get_feature_names_out()
                    self.is_fitted = True
                else:
                    # Transform using existing vocabulary
                    try:
                        tfidf_matrix = self.vectorizer.transform([text])
                    except ValueError as e:
                        # If features mismatch, refit the vectorizer
                        logging.warning(f"Feature mismatch, refitting vectorizer: {e}")
                        self.vectorizer = TfidfVectorizer(
                            max_features=5000,
                            stop_words='english',
                            ngram_range=(1, 3),
                            lowercase=True,
                            strip_accents='unicode'
                        )
                        tfidf_matrix = self.vectorizer.fit_transform([text])
                        self.feature_names = self.vectorizer.get_feature_names_out()
                
                # Get top topics
                topics = []
                feature_names = self.vectorizer.get_feature_names_out()
                
                if tfidf_matrix.shape[1] > 0:  # Check if we have any features
                    scores = tfidf_matrix.toarray()[0]
                    sorted_idx = np.argsort(scores)[::-1]
                    
                    # Get top 10 topics with non-zero scores
                    for idx in sorted_idx[:10]:
                        if scores[idx] > 0:
                            topics.append({
                                'word': feature_names[idx],
                                'frequency': float(scores[idx]),
                                'importance': float(scores[idx])
                            })
                
                return topics
                
        except Exception as e:
            logging.error(f"Error extracting topics: {e}")
            return []

    def analyze_content_parallel(self, file_path, content, metadata=None):
        """Analyze content using both GPU and CPU in parallel."""
        try:
            # Generate content hash for deduplication
            content_hash = hashlib.md5(content.encode()).hexdigest()
            
            # Extract basic metadata
            file_info = {
                'path': file_path,
                'name': os.path.basename(file_path),
                'extension': os.path.splitext(file_path)[1].lower(),
                'size': os.path.getsize(file_path),
                'modified': datetime.fromtimestamp(os.path.getmtime(file_path)),
                'content_hash': content_hash,
                'content_length': len(content),
                'word_count': len(content.split()),
                'metadata': metadata or {}
            }
            
            # Split content into chunks for parallel processing
            chunks = self._split_content(content)
            
            # Initialize results
            topics = []
            entities = []
            vectors = []
            
            # Process chunks with proper thread pool management
            with ThreadPoolExecutor(max_workers=2) as executor:
                try:
                    futures = []
                    
                    # Submit GPU tasks if available
                    if self.gpu_model is not None:
                        for chunk in chunks:
                            if chunk.strip():  # Only process non-empty chunks
                                futures.append(executor.submit(
                                    self._process_chunk_gpu,
                                    chunk,
                                    file_path
                                ))
                    
                    # Submit CPU tasks
                    if self.cpu_model is not None:
                        for chunk in chunks:
                            if chunk.strip():  # Only process non-empty chunks
                                futures.append(executor.submit(
                                    self._process_chunk_cpu,
                                    chunk,
                                    file_path
                                ))
                    
                    # Process results as they complete
                    for future in futures:
                        try:
                            result = future.result(timeout=30)
                            if result is not None:
                                chunk_topics, chunk_entities, chunk_vector = result
                                if chunk_topics:
                                    topics.extend(chunk_topics)
                                if chunk_entities:
                                    entities.extend(chunk_entities)
                                if chunk_vector is not None and chunk_vector.size > 0:
                                    vectors.append(chunk_vector)
                        except Exception as e:
                            logging.error(f"Error processing chunk: {e}")
                            continue
                            
                except Exception as e:
                    logging.error(f"Error in thread pool execution: {e}")
                finally:
                    # Ensure proper cleanup
                    executor.shutdown(wait=True)
            
            # Combine results
            file_info['topics'] = self._combine_topics(topics) if topics else []
            file_info['entities'] = self._combine_entities(entities) if entities else []
            
            # Calculate final content vector
            if vectors:
                try:
                    # Ensure all vectors have the same shape
                    max_features = max(v.shape[0] for v in vectors)
                    padded_vectors = []
                    for v in vectors:
                        if v.shape[0] < max_features:
                            padded = np.zeros(max_features)
                            padded[:v.shape[0]] = v
                            padded_vectors.append(padded)
                        else:
                            padded_vectors.append(v)
                    if padded_vectors:
                        self.content_vectors[file_path] = np.mean(padded_vectors, axis=0)
                except Exception as e:
                    logging.error(f"Error calculating content vector: {e}")
            
            # Update file metadata
            self.file_metadata[file_path] = file_info
            
            # Update content graph
            try:
                self._update_content_graph(file_path, file_info.get('topics', []), file_info.get('entities', []))
            except Exception as e:
                logging.error(f"Error updating content graph: {e}")
            
            return file_info
            
        except Exception as e:
            logging.error(f"Error analyzing content for {file_path}: {e}")
            return None

    def _process_chunk_gpu(self, chunk, file_path):
        """Process a chunk using GPU."""
        try:
            if not self.gpu_model or not chunk.strip():
                return None
                
            # Extract topics with error handling
            try:
                topics = self._extract_topics(chunk)
            except Exception as e:
                logging.error(f"Error extracting topics in GPU processing: {e}")
                topics = []
            
            # Extract entities with error handling
            try:
                entities = self._extract_entities(chunk)
            except Exception as e:
                logging.error(f"Error extracting entities in GPU processing: {e}")
                entities = []
            
            # Calculate vector with error handling
            try:
                with self.lock:
                    self._initialize_vectorizer()
                    if not self.is_fitted:
                        vector = self.vectorizer.fit_transform([chunk]).toarray()[0]
                        self.feature_names = self.vectorizer.get_feature_names_out()
                        self.is_fitted = True
                    else:
                        vector = self.vectorizer.transform([chunk]).toarray()[0]
            except Exception as e:
                logging.error(f"Error calculating vector in GPU processing: {e}")
                vector = None
            
            return topics, entities, vector
            
        except Exception as e:
            logging.error(f"Error processing chunk on GPU: {e}")
            return None

    def _process_chunk_cpu(self, chunk, file_path):
        """Process a chunk using CPU."""
        try:
            if not self.cpu_model or not chunk.strip():
                return None
                
            # Extract topics with error handling
            try:
                topics = self._extract_topics(chunk)
            except Exception as e:
                logging.error(f"Error extracting topics in CPU processing: {e}")
                topics = []
            
            # Extract entities with error handling
            try:
                entities = self._extract_entities(chunk)
            except Exception as e:
                logging.error(f"Error extracting entities in CPU processing: {e}")
                entities = []
            
            # Calculate vector with error handling
            try:
                with self.lock:
                    self._initialize_vectorizer()
                    if not self.is_fitted:
                        vector = self.vectorizer.fit_transform([chunk]).toarray()[0]
                        self.feature_names = self.vectorizer.get_feature_names_out()
                        self.is_fitted = True
                    else:
                        vector = self.vectorizer.transform([chunk]).toarray()[0]
            except Exception as e:
                logging.error(f"Error calculating vector in CPU processing: {e}")
                vector = None
            
            return topics, entities, vector
            
        except Exception as e:
            logging.error(f"Error processing chunk on CPU: {e}")
            return None

    def _split_content(self, content, chunk_size=512):
        """Split content into chunks for parallel processing."""
        words = content.split()
        chunks = []
        
        for i in range(0, len(words), chunk_size):
            chunk = ' '.join(words[i:i + chunk_size])
            chunks.append(chunk)
            
        return chunks
        
    def _combine_topics(self, topics_list):
        """Combine topics from multiple chunks."""
        try:
            # Combine all topics
            all_topics = []
            for topics in topics_list:
                if isinstance(topics, list):
                    all_topics.extend(topics)
            
            # Calculate combined frequencies
            topic_freq = defaultdict(float)
            for topic in all_topics:
                if isinstance(topic, dict) and 'word' in topic and 'frequency' in topic:
                    topic_freq[topic['word']] += float(topic['frequency'])
            
            # Create combined topics
            combined_topics = []
            total_freq = sum(topic_freq.values()) or 1.0  # Avoid division by zero
            
            for word, freq in sorted(topic_freq.items(), key=lambda x: x[1], reverse=True)[:10]:
                combined_topics.append({
                    'word': word,
                    'frequency': freq,
                    'importance': freq / total_freq
                })
            
            return combined_topics
            
        except Exception as e:
            logging.error(f"Error combining topics: {e}")
            return []
            
    def _combine_entities(self, entities_list):
        """Combine entities from multiple chunks."""
        try:
            # Combine all entities
            all_entities = []
            for entities in entities_list:
                if isinstance(entities, list):
                    all_entities.extend(entities)
            
            # Remove duplicates while preserving order
            seen = set()
            unique_entities = []
            for entity in all_entities:
                if isinstance(entity, dict) and 'text' in entity:
                    if entity['text'] not in seen:
                        seen.add(entity['text'])
                        unique_entities.append(entity)
            
            return unique_entities
            
        except Exception as e:
            logging.error(f"Error combining entities: {e}")
            return []

    def _extract_entities(self, text):
        """Extract named entities from text using regex patterns."""
        try:
            # Clean and preprocess text
            text = clean_text(text)
            if not text:
                return []

            entities = []
            
            # Define regex patterns for different entity types
            patterns = {
                'person': r'\b[A-Z][a-z]+(?:\s+[A-Z][a-z]+)*\b',  # Names
                'organization': r'\b[A-Z][A-Za-z]+(?:\s+[A-Z][A-Za-z]+)*\b',  # Organizations
                'date': r'\b\d{1,2}[-/]\d{1,2}[-/]\d{2,4}\b|\b(?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)[a-z]* \d{1,2},? \d{4}\b',  # Dates
                'email': r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b',  # Email addresses
                'url': r'https?://(?:[-\w.]|(?:%[\da-fA-F]{2}))+',  # URLs
                'number': r'\b\d+(?:,\d+)*(?:\.\d+)?\b',  # Numbers
                'code': r'`[^`]+`|```[\s\S]+?```'  # Code blocks
            }
            
            # Extract entities for each pattern
            for entity_type, pattern in patterns.items():
                matches = re.finditer(pattern, text)
                for match in matches:
                    entity_text = match.group()
                    # Skip very short matches and common words
                    if len(entity_text) > 2 and not entity_text.lower() in {'the', 'and', 'for', 'but', 'nor', 'yet'}:
                        entities.append({
                            'text': entity_text,
                            'type': entity_type,
                            'start': match.start(),
                            'end': match.end()
                        })
            
            # Remove duplicates while preserving order
            seen = set()
            unique_entities = []
            for entity in entities:
                if entity['text'] not in seen:
                    seen.add(entity['text'])
                    unique_entities.append(entity)
            
            return unique_entities
            
        except Exception as e:
            logging.error(f"Error extracting entities: {e}")
            return []

    def _update_content_graph(self, file_path, topics, entities):
        """Update the content relationship graph with new file information."""
        try:
            # Add file node if not exists
            if file_path not in self.content_graph:
                self.content_graph.add_node(file_path, type='file')
            
            # Add topic nodes and edges
            for topic in topics:
                topic_word = topic['word']
                if topic_word not in self.content_graph:
                    self.content_graph.add_node(topic_word, type='topic')
                self.content_graph.add_edge(file_path, topic_word, 
                                         weight=topic['importance'],
                                         type='has_topic')
            
            # Add entity nodes and edges
            for entity in entities:
                entity_text = entity['text']
                if entity_text not in self.content_graph:
                    self.content_graph.add_node(entity_text, type='entity')
                self.content_graph.add_edge(file_path, entity_text,
                                         weight=1.0,
                                         type='has_entity')
            
            # Update relationship cache
            self.relationship_cache[file_path] = {
                'topics': topics,
                'entities': entities
            }
            
        except Exception as e:
            logging.error(f"Error updating content graph: {e}")

    def get_content_context(self, file_path):
        """Get the context information for a file."""
        try:
            if file_path not in self.file_metadata:
                return None
                
            file_info = self.file_metadata[file_path]
            
            # Get related files
            related_files = []
            if file_path in self.content_graph:
                for neighbor in self.content_graph.neighbors(file_path):
                    if self.content_graph.nodes[neighbor]['type'] == 'file':
                        edge_data = self.content_graph.get_edge_data(file_path, neighbor)
                        related_files.append({
                            'path': neighbor,
                            'relationship_type': edge_data.get('type', 'unknown'),
                            'similarity': edge_data.get('weight', 0.0)
                        })
            
            # Get common topics with related files
            common_topics = []
            if related_files:
                for rel_file in related_files:
                    if rel_file['path'] in self.relationship_cache:
                        file_topics = set(t['word'] for t in self.relationship_cache[file_path]['topics'])
                        rel_topics = set(t['word'] for t in self.relationship_cache[rel_file['path']]['topics'])
                        common = file_topics.intersection(rel_topics)
                        if common:
                            common_topics.extend(list(common))
            
            # Compile context information
            context = {
                'file_info': {
                    'path': file_path,
                    'name': os.path.basename(file_path),
                    'extension': os.path.splitext(file_path)[1].lower(),
                    'size': file_info.get('size', 0),
                    'modified': file_info.get('modified', None),
                    'content_length': file_info.get('content_length', 0),
                    'word_count': file_info.get('word_count', 0)
                },
                'content_summary': {
                    'topics': file_info.get('topics', []),
                    'entities': file_info.get('entities', [])
                },
                'relationships': {
                    'related_files': related_files,
                    'common_topics': list(set(common_topics))
                }
            }
            
            return context
            
        except Exception as e:
            logging.error(f"Error getting content context for {file_path}: {e}")
            return None

# Initialize the content analyzer
content_analyzer = ContentAnalyzer()

def collect_all_files(directory):
    """Recursively collect all files from a directory."""
    all_files = []
    try:
        # Skip these directories
        skip_dirs = {'.git', '__pycache__', 'node_modules', '.vscode', '.idea'}
        
        # Binary and media file extensions to skip content analysis
        binary_extensions = {
            '.mp3', '.mp4', '.wav', '.avi', '.mov', '.flv', '.wmv',  # Audio/Video
            '.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp',  # Images
            '.zip', '.rar', '.7z', '.tar', '.gz', '.exe', '.dll',  # Binary/Archives
            '.db', '.dat', '.bin', '.pyc', '.pyo', '.pyd'  # Data files
        }
        
        for root, dirs, files in os.walk(directory):
            # Remove skip directories
            dirs[:] = [d for d in dirs if d not in skip_dirs and not d.startswith('.')]
            
            for file in files:
                try:
                    file_path = os.path.join(root, file)
                    ext = os.path.splitext(file)[1].lower()
                    
                    # Skip binary and media files from content analysis
                    if ext in binary_extensions:
                        if log_callback:
                            log_callback(f"Skipping binary file: {file}")
                        continue
                        
                    # Skip if not a regular file or empty
                    if not os.path.isfile(file_path) or os.path.getsize(file_path) == 0:
                        continue
                        
                    all_files.append(file_path)
                    
                except Exception as e:
                    logging.error(f"Error accessing file {file} in {root}: {e}")
                    continue
            
        return all_files
        
    except Exception as e:
        logging.error(f"Error walking directory {directory}: {e}")
        return all_files

def _determine_primary_topic(file_info, context):
    """Determine the primary topic based on file info and context."""
    try:
        # Get topics from file info
        topics = file_info.get('topics', [])
        if not topics:
            return "Uncategorized"
            
        # Get the most important topic
        primary_topic = max(topics, key=lambda x: x.get('importance', 0))
        return primary_topic['word'].title()
        
    except Exception as e:
        logging.error(f"Error determining primary topic: {e}")
        return "Uncategorized"

def _determine_topic_hierarchy(file_info, context):
    """Build a topic hierarchy based on content analysis."""
    try:
        # Get topics and entities
        topics = file_info.get('topics', [])
        entities = file_info.get('entities', [])
        
        # Build hierarchy levels
        hierarchy = []
        
        # Level 1: Main topic (from primary topic)
        main_topic = _determine_primary_topic(file_info, context)
        hierarchy.append(main_topic)
        
        # Level 2: Subtopics (from remaining important topics)
        subtopics = []
        for topic in topics[1:3]:  # Take next 2 most important topics
            if topic['importance'] > 0.3:  # Only include significant topics
                subtopics.append(topic['word'].title())
        if subtopics:
            hierarchy.append(subtopics)
        
        # Level 3: Context-based categorization
        if context and 'domain_indicators' in context:
            domain = list(context['domain_indicators'])[0]
            hierarchy.append(domain)
        
        return hierarchy
        
    except Exception as e:
        logging.error(f"Error determining topic hierarchy: {e}")
        return ["Uncategorized"]

def _determine_relationships(file_info, context):
    """Determine relationships between files based on content and context."""
    try:
        relationships = {
            'similar_content': [],
            'referenced_by': [],
            'references_to': [],
            'temporal_related': []
        }
        
        if not context:
            return relationships
            
        # Add similar content files
        if 'similar_files' in context:
            relationships['similar_content'] = [
                {'file': f['path'], 'similarity': f['similarity']}
                for f in context['similar_files']
            ]
        
        # Add referenced files
        if 'references' in context:
            relationships['references_to'] = [
                {'file': ref['path'], 'type': ref['type']}
                for ref in context['references']
            ]
        
        # Add temporal relationships
        if 'temporal_markers' in context:
            relationships['temporal_related'] = [
                {'file': tm['path'], 'relation': tm['relation']}
                for tm in context.get('temporal_related', [])
            ]
        
        return relationships
        
    except Exception as e:
        logging.error(f"Error determining relationships: {e}")
        return {}

def _determine_content_context(file_info, context):
    """Determine content context including domain, temporal markers, and relationships."""
    try:
        content_context = {
            'domain_indicators': set(),
            'temporal_markers': [],
            'main_topics': [],
            'relationships': []
        }
        
        # Add domain indicators
        for domain, info in DOMAIN_INDICATORS.items():
            patterns = info['patterns']
            content = file_info.get('content', '')
            if any(pattern in content.lower() for pattern in patterns):
                content_context['domain_indicators'].add(domain)
        
        # Add temporal markers
        if context and 'temporal_markers' in context:
            content_context['temporal_markers'] = context['temporal_markers']
        
        # Add main topics
        if file_info.get('topics'):
            content_context['main_topics'] = sorted(
                file_info['topics'],
                key=lambda x: x.get('importance', 0),
                reverse=True
            )[:5]  # Top 5 topics
        
        # Add relationships
        if context and 'relationships' in context:
            content_context['relationships'] = context['relationships']
        
        return content_context
        
    except Exception as e:
        logging.error(f"Error determining content context: {e}")
        return {}

def get_folder_structure(directory, prefix="", show_details=False):
    """Generate a visual representation of the folder structure with optional details."""
    try:
        if not os.path.exists(directory):
            return "Directory does not exist"
            
        structure = []
        
        # Add root directory name
        if not prefix:
            structure.append(os.path.basename(directory) + "/")
        
        # Get all items in directory
        items = os.listdir(directory)
        items.sort()
        
        # Process each item
        for i, item in enumerate(items):
            path = os.path.join(directory, item)
            is_last = i == len(items) - 1
            
            # Create branch symbol
            branch = "└── " if is_last else "├── "
            
            # Add item details
            if show_details and os.path.isfile(path):
                size = os.path.getsize(path)
                modified = datetime.fromtimestamp(os.path.getmtime(path))
                details = f" ({size/1024:.1f}KB, modified: {modified.strftime('%Y-%m-%d %H:%M')})"
            else:
                details = "/" if os.path.isdir(path) else ""
            
            # Add item to structure with proper encoding
            line = f"{prefix}{branch}{item}{details}"
            structure.append(line.encode('utf-8', errors='replace').decode('utf-8'))
            
            # If directory, recursively process it
            if os.path.isdir(path):
                # Create extended prefix for subdirectories
                ext_prefix = prefix + ("    " if is_last else "│   ")
                substructure = get_folder_structure(path, ext_prefix, show_details)
                if substructure:
                    structure.extend(substructure.split('\n'))
        
        return '\n'.join(structure)
        
    except Exception as e:
        logging.error(f"Error getting folder structure: {e}")
        return "Error generating folder structure"

def determine_project_structure(content, context):
    """Determine project name and structure based on content and context."""
    try:
        # Extract project indicators
        project_indicators = {
            'Web Application': ['web', 'app', 'react', 'angular', 'vue', 'html', 'css', 'javascript'],
            'Mobile App': ['android', 'ios', 'mobile', 'app', 'flutter', 'react native'],
            'Data Science': ['data', 'analysis', 'machine learning', 'model', 'dataset'],
            'Backend Service': ['api', 'server', 'database', 'service', 'microservice'],
            'Documentation': ['docs', 'guide', 'manual', 'readme', 'wiki'],
            'Development Tools': ['tool', 'utility', 'cli', 'script', 'automation']
        }

        # Standard project structure
        standard_structure = {
            'Documentation': {
                'Technical': ['api', 'architecture', 'database'],
                'User': ['guides', 'manuals', 'tutorials'],
                'Project': ['requirements', 'planning', 'meetings']
            },
            'Source': {
                'Main': ['src', 'app', 'core'],
                'Tests': ['unit', 'integration', 'e2e'],
                'Scripts': ['build', 'deploy', 'utilities']
            },
            'Resources': {
                'Assets': ['images', 'fonts', 'media'],
                'Data': ['configs', 'samples', 'schemas'],
                'External': ['libraries', 'dependencies']
            },
            'Build': {
                'Output': ['dist', 'build', 'release'],
                'Temp': ['cache', 'temp', 'logs']
            }
        }

        # Determine project type
        project_type = 'General Project'
        content_lower = content.lower()
        for proj_type, indicators in project_indicators.items():
            if any(ind in content_lower for ind in indicators):
                project_type = proj_type
                break

        # Get project name from context
        project_name = None
        if context.get('metadata', {}).get('filename'):
            name_parts = os.path.splitext(context['metadata']['filename'])[0].split('_')
            project_name = ' '.join(word.title() for word in name_parts if len(word) > 2)
        
        if not project_name and context.get('topics'):
            project_name = f"{project_type}_{context['topics'][0].title()}"
        
        project_name = project_name or project_type

        # Customize structure based on project type
        structure = {
            'name': project_name,
            'type': project_type,
            'folders': {}
        }

        # Add standard structure
        structure['folders'] = standard_structure.copy()

        # Add project-specific folders
        if project_type == 'Web Application':
            structure['folders'].update({
                'Frontend': {
                    'Components': ['ui', 'pages', 'layouts'],
                    'Assets': ['styles', 'images', 'fonts'],
                    'State': ['store', 'context', 'reducers']
                },
                'Backend': {
                    'API': ['routes', 'controllers', 'middleware'],
                    'Database': ['models', 'migrations', 'seeds'],
                    'Services': ['auth', 'email', 'storage']
                }
            })
        elif project_type == 'Data Science':
            structure['folders'].update({
                'Data': {
                    'Raw': ['datasets', 'sources', 'external'],
                    'Processed': ['cleaned', 'transformed', 'features'],
                    'Models': ['trained', 'evaluation', 'predictions']
                },
                'Notebooks': {
                    'Analysis': ['eda', 'visualization', 'reports'],
                    'Models': ['training', 'testing', 'validation'],
                    'Research': ['experiments', 'literature', 'results']
                }
            })
        elif project_type == 'Backend Service':
            structure['folders'].update({
                'API': {
                    'Routes': ['v1', 'v2', 'public'],
                    'Controllers': ['auth', 'data', 'admin'],
                    'Middleware': ['validation', 'security', 'logging']
                },
                'Infrastructure': {
                    'Database': ['migrations', 'seeds', 'backup'],
                    'Services': ['queue', 'cache', 'storage'],
                    'Config': ['env', 'secrets', 'settings']
                }
            })

        return structure

    except Exception as e:
        logging.error(f"Error determining project structure: {e}")
        return {
            'name': 'General Project',
            'type': 'General',
            'folders': standard_structure
        }

def organize_by_project_structure(analyses, dest_root, progress_callback=None, cancel_event=None, log_callback=None):
    """Organize files into project-based structures."""
    try:
        # Group files by potential projects
        projects = {}
        
        # First pass: Identify projects and their files
        for file_path, analysis in analyses.items():
            content = analysis['content']
            context = analysis['context']
            
            # Get project structure
            structure = determine_project_structure(content, context)
            project_name = structure['name']
            
            if project_name not in projects:
                projects[project_name] = {
                    'structure': structure,
                    'files': [],
                    'context': context
                }
            projects[project_name]['files'].append((file_path, analysis))

        # Second pass: Organize files within each project
        total_files = sum(len(project['files']) for project in projects.values())
        organized = 0

        for project_name, project_info in projects.items():
            if cancel_event and cancel_event.is_set():
                return

            try:
                structure = project_info['structure']
                project_root = os.path.join(dest_root, project_name)

                # Create project directory structure
                for main_folder, subfolders in structure['folders'].items():
                    for subfolder, categories in subfolders.items():
                        for category in categories:
                            folder_path = os.path.join(project_root, main_folder, subfolder, category)
                            os.makedirs(folder_path, exist_ok=True)

                # Organize files into appropriate folders
                for file_path, analysis in project_info['files']:
                    if not os.path.exists(file_path):
                        continue

                    content = analysis['content']
                    context = analysis['context']

                    # Determine appropriate subfolder
                    main_folder, subfolder = determine_file_location(content, context)
                    
                    # Create destination path
                    dest_folder = os.path.join(project_root, main_folder, subfolder)
                    os.makedirs(dest_folder, exist_ok=True)
                    
                    # Move file
                    dest_path = os.path.join(dest_folder, os.path.basename(file_path))
                    shutil.move(file_path, dest_path)
                    
                    organized += 1
                    if progress_callback:
                        progress = int((organized / total_files) * 100)
                        progress_callback(progress, f"Organizing {project_name}...")

                if log_callback:
                    log_callback(f"\nCreated project structure for: {project_name}")
            except Exception as e:
                logging.error(f"Error organizing project {project_name}: {e}")

        # Clean up empty folders
        clean_empty_folders(dest_root)

    except Exception as e:
        logging.error(f"Error in project organization: {e}")
        if log_callback:
            log_callback(f"Error: {e}")

def determine_file_location(content, context):
    """Determine the appropriate location for a file within the project structure."""
    try:
        content_lower = content.lower()
        
        # Check for documentation
        if any(term in content_lower for term in ['readme', 'guide', 'manual', 'doc']):
            if 'api' in content_lower or 'architecture' in content_lower:
                return 'Documentation', 'Technical'
            elif 'user' in content_lower or 'tutorial' in content_lower:
                return 'Documentation', 'User'
            return 'Documentation', 'Project'

        # Check for source code
        if context.get('content_features', {}).get('has_code'):
            if 'test' in content_lower:
                return 'Source', 'Tests'
            elif 'script' in content_lower or 'tool' in content_lower:
                return 'Source', 'Scripts'
            return 'Source', 'Main'

        # Check for resources
        if any(term in content_lower for term in ['image', 'font', 'media', 'asset']):
            return 'Resources', 'Assets'
        elif any(term in content_lower for term in ['config', 'setting', 'env']):
            return 'Resources', 'Data'
        elif any(term in content_lower for term in ['lib', 'vendor', 'package']):
            return 'Resources', 'External'

        # Check for build artifacts
        if any(term in content_lower for term in ['build', 'dist', 'release']):
            return 'Build', 'Output'
        elif any(term in content_lower for term in ['temp', 'cache', 'log']):
            return 'Build', 'Temp'

        # Default location
        return 'Source', 'Main'

    except Exception as e:
        logging.error(f"Error determining file location: {e}")
        return 'Source', 'Main'

def organize_by_content_similarity(src_root, dest_root, progress_callback=None, cancel_event=None, log_callback=None):
    """Organize files based on content similarity with parallel processing."""
    if progress_callback:
        progress_callback(0, "Starting content analysis...")
    
    try:
        # Print original structure
        if log_callback:
            log_callback("\nOriginal File Structure:")
            log_callback("=" * 50)
            structure = get_folder_structure(src_root, show_details=True)
            for line in structure.split('\n'):
                log_callback(line)
            log_callback("=" * 50 + "\n")
        
        # Collect all files
        files = collect_all_files(src_root)
        if not files:
            if log_callback:
                log_callback("No files found to analyze")
            return
        
        if log_callback:
            log_callback(f"\nFound {len(files)} files to process")
        
        # First phase: Batch analyze all files
        analyses, binary_files = batch_analyze_files(
            files,
            progress_callback,
            log_callback
        )
        
        if cancel_event and cancel_event.is_set():
            return
            
        # Organize files into project structures
        organize_by_project_structure(analyses, dest_root, progress_callback, cancel_event, log_callback)
        
        # Handle binary files
        if binary_files:
            organize_binary_files(binary_files, dest_root, progress_callback)
        
        # Print new structure at the end
        if log_callback:
            log_callback("\nNew File Structure:")
            log_callback("=" * 50)
            structure = get_folder_structure(dest_root, show_details=True)
            for line in structure.split('\n'):
                log_callback(line)
            log_callback("=" * 50 + "\n")
            
            # Print summary
            log_callback("\n=== Organization Summary ===")
            category_counts = defaultdict(int)
            for category, files in analyses.items():
                category_counts[category] = len(files)
                log_callback(f"{category}: {category_counts[category]} files")
        
        if progress_callback:
            progress_callback(100, "Organization completed")
            
    except Exception as e:
        logging.error(f"Error in organize_by_content_similarity: {e}")
        if log_callback:
            log_callback(f"Error: {str(e)}")

def organize_binary_files(binary_files, dest_root, progress_callback=None):
    """Organize binary files into appropriate folders."""
    try:
        total_files = len(binary_files)
        organized = 0
        
        for file_path in binary_files:
            ext = os.path.splitext(file_path)[1].lower()
            
            # Determine category and project folder
            if ext in {'.jpg', '.jpeg', '.png', '.gif'}:
                category = 'Resources/Assets/Images'
            elif ext in {'.mp3', '.wav'}:
                category = 'Resources/Assets/Audio'
            elif ext in {'.mp4', '.avi', '.mov'}:
                category = 'Resources/Assets/Video'
            elif ext in {'.pdf', '.doc', '.docx'}:
                category = 'Documentation/External'
            elif ext in {'.zip', '.rar', '.7z'}:
                category = 'Resources/External/Archives'
            else:
                category = 'Resources/Other'
            
            # Create destination directory
            dest_dir = os.path.join(dest_root, category)
            os.makedirs(dest_dir, exist_ok=True)
            
            # Move file
            dest_path = os.path.join(dest_dir, os.path.basename(file_path))
            shutil.move(file_path, dest_path)
            
            organized += 1
            if progress_callback:
                progress = int((organized / total_files) * 100)
                progress_callback(progress, "Organizing media files...")
                
    except Exception as e:
        logging.error(f"Error organizing binary files: {e}")

class SmartContentAnalyzer:
    def __init__(self):
        # Domain-specific keywords for quick content analysis
        self.domain_keywords = {
            'Web Development': {
                'keywords': {'html', 'css', 'javascript', 'react', 'angular', 'vue', 'frontend', 'backend', 'api', 'web', 'http', 'rest', 'server'},
                'priority': 10,
                'folders': ['Frontend', 'Backend', 'Assets', 'API']
            },
            'Data Science': {
                'keywords': {'data', 'analysis', 'model', 'train', 'predict', 'dataset', 'pandas', 'numpy', 'sklearn', 'tensorflow', 'pytorch'},
                'priority': 10,
                'folders': ['Data', 'Models', 'Analysis', 'Notebooks']
            },
            'Documentation': {
                'keywords': {'readme', 'documentation', 'guide', 'manual', 'tutorial', 'howto', 'setup', 'install'},
                'priority': 8,
                'folders': ['Guides', 'API', 'Technical', 'User']
            },
            'Development': {
                'keywords': {'src', 'test', 'lib', 'package', 'module', 'class', 'function', 'method', 'implementation'},
                'priority': 7,
                'folders': ['Source', 'Tests', 'Libraries', 'Utils']
            }
        }

        # File type categories
        self.file_categories = {
            'source_code': {
                'extensions': {'.py', '.java', '.js', '.cpp', '.h', '.cs', '.php', '.rb', '.go'},
                'folders': ['Source', 'Implementation']
            },
            'web': {
                'extensions': {'.html', '.css', '.js', '.jsx', '.ts', '.tsx', '.vue', '.php'},
                'folders': ['Frontend', 'Web']
            },
            'data': {
                'extensions': {'.csv', '.json', '.xml', '.yaml', '.sql', '.db'},
                'folders': ['Data', 'Resources']
            },
            'documentation': {
                'extensions': {'.md', '.txt', '.pdf', '.doc', '.docx'},
                'folders': ['Documentation', 'Docs']
            },
            'config': {
                'extensions': {'.config', '.ini', '.env', '.yml', '.yaml', '.json'},
                'folders': ['Config', 'Settings']
            }
        }

    def analyze_file(self, file_path):
        """Analyze a file using smart keyword extraction and context analysis."""
        try:
            # Get file metadata
            filename = os.path.basename(file_path)
            ext = os.path.splitext(file_path)[1].lower()
            stats = os.stat(file_path)

            # Initialize analysis results
            analysis = {
                'path': file_path,
                'filename': filename,
                'extension': ext,
                'size': stats.st_size,
                'modified': stats.st_mtime,
                'domains': [],
                'keywords': set(),
                'category': None,
                'context': {},
                'suggested_folders': []
            }

            # Skip binary files
            if self._is_binary_file(file_path, ext):
                analysis['category'] = 'binary'
                analysis['suggested_folders'] = ['Resources', 'Binary']
                return analysis

            # Read file content
            content = self._read_file_content(file_path)
            if not content:
                return analysis

            # Extract keywords and determine domains
            keywords = self._extract_keywords(content)
            analysis['keywords'] = keywords

            # Determine domains and context
            domains = self._determine_domains(keywords, ext)
            analysis['domains'] = domains

            # Determine file category
            category = self._determine_category(ext, keywords)
            analysis['category'] = category

            # Build context based on content and metadata
            context = self._build_context(content, keywords, domains, category)
            analysis['context'] = context

            # Suggest folders based on analysis
            folders = self._suggest_folders(domains, category, context)
            analysis['suggested_folders'] = folders

            return analysis

        except Exception as e:
            logging.error(f"Error analyzing file {file_path}: {e}")
            return None

    def _is_binary_file(self, file_path, ext):
        """Check if file is binary."""
        binary_extensions = {
            '.exe', '.dll', '.so', '.dylib', '.bin',
            '.jpg', '.jpeg', '.png', '.gif', '.bmp',
            '.mp3', '.mp4', '.wav', '.avi', '.mov',
            '.zip', '.rar', '.7z', '.tar', '.gz'
        }
        return ext in binary_extensions

    def _read_file_content(self, file_path, max_size=1024*1024):  # 1MB limit
        """Read file content with size limit."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read(max_size)
        except Exception as e:
            logging.warning(f"Error reading file {file_path}: {e}")
            return ""

    def _extract_keywords(self, content):
        """Extract important keywords from content."""
        # Convert to lowercase and split into words
        words = set(word.lower() for word in content.split())
        
        # Remove common words and short terms
        stop_words = {'the', 'a', 'an', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for', 'of', 'with', 'by'}
        keywords = {word for word in words 
                   if len(word) > 2 
                   and word not in stop_words 
                   and not word.isnumeric()}
        
        return keywords

    def _determine_domains(self, keywords, ext):
        """Determine relevant domains based on keywords and extension."""
        domains = []
        for domain, info in self.domain_keywords.items():
            # Calculate match score based on keyword matches
            matches = keywords.intersection(info['keywords'])
            if matches:
                score = len(matches) * info['priority']
                domains.append((domain, score))
        
        # Sort domains by score and return top 2
        domains.sort(key=lambda x: x[1], reverse=True)
        return [domain for domain, _ in domains[:2]]

    def _determine_category(self, ext, keywords):
        """Determine file category based on extension and keywords."""
        for category, info in self.file_categories.items():
            if ext in info['extensions']:
                return category
        return 'other'

    def _build_context(self, content, keywords, domains, category):
        """Build context information for the file."""
        context = {
            'primary_domain': domains[0] if domains else None,
            'secondary_domain': domains[1] if len(domains) > 1 else None,
            'category': category,
            'key_terms': list(keywords)[:10],  # Top 10 keywords
            'content_type': self._guess_content_type(content, keywords)
        }
        return context

    def _guess_content_type(self, content, keywords):
        """Guess the type of content based on its structure and keywords."""
        content_lower = content.lower()
        
        if any(word in content_lower for word in {'class', 'function', 'def', 'import', 'return'}):
            return 'source_code'
        elif any(word in content_lower for word in {'test', 'assert', 'describe', 'it', 'should'}):
            return 'test'
        elif any(word in content_lower for word in {'select', 'insert', 'update', 'delete', 'from', 'where'}):
            return 'database'
        elif any(word in content_lower for word in {'<html>', '<body>', '<div>', '<script>'}):
            return 'web'
        elif len(content.splitlines()) < 50 and any(word in content_lower for word in {'config', 'setting', 'env', 'variable'}):
            return 'configuration'
        else:
            return 'document'

    def _suggest_folders(self, domains, category, context):
        """Suggest folder structure based on analysis."""
        folders = []
        
        # Add domain-based folders
        if domains:
            primary_domain = domains[0]
            if primary_domain in self.domain_keywords:
                folders.extend(self.domain_keywords[primary_domain]['folders'][:2])
        
        # Add category-based folders
        if category in self.file_categories:
            folders.extend(self.file_categories[category]['folders'][:2])
        
        # Add context-based folder
        content_type = context.get('content_type')
        if content_type == 'test':
            folders.append('Tests')
        elif content_type == 'configuration':
            folders.append('Config')
        elif content_type == 'database':
            folders.append('Database')
        
        # Ensure we have at least one folder
        if not folders:
            folders = ['Other']
        
        return folders

# Initialize the smart analyzer
smart_analyzer = SmartContentAnalyzer()

def smart_organize_files(src_root, dest_root, progress_callback=None, cancel_event=None, log_callback=None):
    """Organize files using smart content analysis."""
    if progress_callback:
        progress_callback(0, "Starting smart analysis...")
    
    try:
        # Print original structure
        if log_callback:
            log_callback("\nOriginal File Structure:")
            log_callback("=" * 50)
            log_callback(get_folder_structure(src_root, show_details=True))
            log_callback("=" * 50 + "\n")
        
        # Collect files
        files = collect_all_files(src_root)
        if not files:
            if log_callback:
                log_callback("No files found to analyze")
            return
        
        total_files = len(files)
        if log_callback:
            log_callback(f"\nFound {total_files} files to process")
        
        # Analyze and organize files
        organized = 0
        analyses = {}
        
        for file_path in files:
            if cancel_event and cancel_event.is_set():
                return
                
            try:
                # Smart analysis
                analysis = smart_analyzer.analyze_file(file_path)
                if not analysis:
                    continue
                
                analyses[file_path] = analysis
                
                # Create folder structure
                folders = analysis['suggested_folders']
                folder_path = os.path.join(dest_root, *folders[:2])  # Use up to 2 levels
                os.makedirs(folder_path, exist_ok=True)
                
                # Move file
                dest_path = os.path.join(folder_path, os.path.basename(file_path))
                shutil.move(file_path, dest_path)
                
                organized += 1
                if progress_callback:
                    progress = int((organized / total_files) * 100)
                    progress_callback(progress, f"Organizing files ({organized}/{total_files})...")
                
                if log_callback:
                    domains = ' & '.join(analysis['domains']) if analysis['domains'] else 'Other'
                    log_callback(f"Organized: {os.path.basename(file_path)} -> {domains}")
                
            except Exception as e:
                logging.error(f"Error processing {file_path}: {e}")
                continue
            
        # Clean up empty folders
        clean_empty_folders(src_root)
        
        # Print new structure
        if log_callback:
            log_callback("\nNew File Structure:")
            log_callback("=" * 50)
            log_callback(get_folder_structure(dest_root, show_details=True))
            log_callback("=" * 50 + "\n")
        
        if progress_callback:
            progress_callback(100, "Smart organization completed")
        
        return analyses
            
    except Exception as e:
        logging.error(f"Error during smart organization: {e}")
        if log_callback:
            log_callback(f"Error: {e}")
        return None

def batch_analyze_files(files, progress_callback=None, log_callback=None):
    """
    Analyze multiple files in parallel with detailed progress reporting.
    
    Args:
        files (list): List of file paths to analyze
        progress_callback (callable): Function to call with progress updates
        log_callback (callable): Function to call with log messages
        
    Returns:
        dict: Analysis results for each file
    """
    if not files:
        return {}
        
    analyzer = SmartContentAnalyzer()
    results = {}
    total_files = len(files)
    processed_files = 0
    error_count = 0
    
    def safe_log(message):
        """Safely call the log callback if provided."""
        if log_callback:
            try:
                log_callback(message)
            except Exception as e:
                logging.error(f"Error in log callback: {e}")
    
    def update_progress():
        """Update progress with detailed statistics."""
        if progress_callback:
            try:
                percentage = (processed_files / total_files) * 100
                status = (f"Analyzed {processed_files}/{total_files} files "
                         f"({percentage:.1f}%) - {error_count} errors")
                progress_callback(percentage, status)
            except Exception as e:
                logging.error(f"Error in progress callback: {e}")
    
    # Group files by type for better progress reporting
    file_types = {}
    for file_path in files:
        ext = os.path.splitext(file_path)[1].lower()
        file_types[ext] = file_types.get(ext, 0) + 1
    
    # Log file type summary
    safe_log("\nStarting batch analysis:")
    for ext, count in file_types.items():
        safe_log(f"- {count} files with extension {ext or '(no extension)'}")
    safe_log("")
    
    for file_path in files:
        try:
            # Get file size for logging
            try:
                file_size = os.path.getsize(file_path)
                size_str = f"{file_size/1024/1024:.1f}MB" if file_size > 1024*1024 else f"{file_size/1024:.1f}KB"
            except Exception as e:
                size_str = "unknown size"
                logging.warning(f"Could not get file size for {file_path}: {e}")
            
            safe_log(f"Analyzing: {os.path.basename(file_path)} ({size_str})")
            
            # Analyze the file
            result = analyzer.analyze_file(file_path)
            
            if result:
                results[file_path] = result
                safe_log(f"✓ Successfully analyzed {os.path.basename(file_path)}")
                if result.get('domains'):
                    safe_log(f"  Detected domains: {', '.join(result['domains'])}")
            else:
                error_count += 1
                safe_log(f"✗ Failed to analyze {os.path.basename(file_path)}")
                
        except Exception as e:
            error_count += 1
            safe_log(f"✗ Error analyzing {os.path.basename(file_path)}: {str(e)}")
            logging.error(f"Error analyzing {file_path}: {str(e)}")
        
        processed_files += 1
        update_progress()
    
    # Final summary
    safe_log("\nAnalysis complete:")
    safe_log(f"- Successfully analyzed: {len(results)} files")
    safe_log(f"- Failed to analyze: {error_count} files")
    
    if error_count > 0:
        safe_log("\nSome files could not be analyzed. This might be due to:")
        safe_log("- Insufficient permissions")
        safe_log("- Corrupted or invalid file formats")
        safe_log("- Files being too large or binary")
    
    return results

def analyze_content_context(content, file_path):
    """Analyze content and build context information."""
    try:
        # Extract metadata
        filename = os.path.basename(file_path)
        ext = os.path.splitext(file_path)[1].lower()
        file_size = os.path.getsize(file_path)
        modified_time = os.path.getmtime(file_path)

        # Clean and analyze content
        cleaned_content = clean_text(content)
        
        # Extract topics using TF-IDF
        vectorizer = TfidfVectorizer(
            max_features=100,
            stop_words='english',
            ngram_range=(1, 2)
        )
        
        try:
            tfidf_matrix = vectorizer.fit_transform([cleaned_content])
            feature_names = vectorizer.get_feature_names_out()
            scores = tfidf_matrix.toarray()[0]
            
            # Get top topics
            topics = []
            for idx in scores.argsort()[::-1][:5]:  # Top 5 topics
                if scores[idx] > 0:
                    topics.append({
                        'word': feature_names[idx],
                        'importance': float(scores[idx])
                    })
        except Exception as e:
            logging.error(f"Error extracting topics: {e}")
            topics = []

        # Build context
        context = {
            'metadata': {
                'filename': filename,
                'extension': ext,
                'size': file_size,
                'modified': datetime.fromtimestamp(modified_time),
            },
            'topics': topics,
            'content_type': _determine_content_type(content, ext),
            'domain_indicators': _find_domain_indicators(cleaned_content)
        }

        return context

    except Exception as e:
        logging.error(f"Error analyzing content context: {e}")
        return {}

def analyze_and_organize_file(file_path, content):
    """Analyze file content and determine organization structure."""
    try:
        # Get content context
        context = analyze_content_context(content, file_path)
        
        # Determine project structure
        project_structure = determine_project_structure(content, context)
        
        # Build analysis result
        analysis = {
            'path': file_path,
            'primary_topic': project_structure['type'],
            'topic_hierarchy': [project_structure['name']] + list(project_structure['folders'].keys()),
            'relationships': _determine_relationships({}, context),
            'content_context': context,
            'project_structure': {
                'project_name': project_structure['name'],
                'project_parts': list(project_structure['folders'].keys()),
                'project_content': _determine_content_path(content, context),
                'project_context': _determine_context_path(content, context)
            }
        }
        
        return analysis

    except Exception as e:
        logging.error(f"Error analyzing file {file_path}: {e}")
        return None

def _determine_content_type(content, ext):
    """Determine the type of content based on extension and content analysis."""
    # Code files
    if ext in {'.py', '.java', '.js', '.cpp', '.h', '.cs', '.php', '.rb', '.go'}:
        return 'source_code'
    # Web files
    elif ext in {'.html', '.css', '.js', '.jsx', '.ts', '.tsx', '.vue', '.php'}:
        return 'web'
    # Data files
    elif ext in {'.csv', '.json', '.xml', '.yaml', '.sql', '.db'}:
        return 'data'
    # Documentation
    elif ext in {'.md', '.txt', '.pdf', '.doc', '.docx'}:
        return 'documentation'
    # Configuration
    elif ext in {'.config', '.ini', '.env', '.yml', '.yaml', '.json'}:
        return 'configuration'
    # Default
    return 'document'

def _find_domain_indicators(content):
    """Find domain-specific indicators in content."""
    indicators = set()
    content_lower = content.lower()
    
    for domain, info in DOMAIN_INDICATORS.items():
        if any(pattern in content_lower for pattern in info['patterns']):
            indicators.add(domain)
    
    return list(indicators)

def _determine_content_path(content, context):
    """Determine the content-based path for organizing the file."""
    try:
        paths = []
        
        # Add primary domain if available
        if context.get('domain_indicators'):
            paths.append(context['domain_indicators'][0])
        
        # Add content type
        content_type = context.get('content_type', 'Other')
        paths.append(content_type.title())
        
        # Add topic-based subfolder
        if context.get('topics'):
            paths.append(context['topics'][0]['word'].title())
        
        return paths if paths else ['Uncategorized']
        
    except Exception as e:
        logging.error(f"Error determining content path: {e}")
        return ['Uncategorized']

def _determine_context_path(content, context):
    """Determine the context-based path for organizing the file."""
    try:
        paths = []
        
        # Add domain context
        if context.get('domain_indicators'):
            paths.append('By Domain')
            paths.append(context['domain_indicators'][0])
        
        # Add temporal context if available
        if context.get('metadata', {}).get('modified'):
            modified = context['metadata']['modified']
            paths.append('By Date')
            paths.append(modified.strftime('%Y-%m'))
        
        return paths if paths else ['Other Context']
        
    except Exception as e:
        logging.error(f"Error determining context path: {e}")
        return ['Other Context']

def clean_empty_folders(root_dir, exclude_dirs=None):
    """
    Recursively remove empty folders from the root directory.
    Args:
        root_dir (str): The root directory to clean.
        exclude_dirs (set): A set of directories to exclude from deletion.
    """
    exclude_dirs = exclude_dirs or set()
    for dirpath, dirnames, _ in os.walk(root_dir, topdown=False):
        for dirname in dirnames:
            folder_path = os.path.join(dirpath, dirname)
            if folder_path not in exclude_dirs and not os.listdir(folder_path):
                try:
                    os.rmdir(folder_path)
                    logging.info(f"Removed empty folder: {folder_path}")
                except Exception as e:
                    logging.error(f"Failed to remove folder {folder_path}: {e}")